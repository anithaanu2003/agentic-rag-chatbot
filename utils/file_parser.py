import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation

def parse_file(file):
    name = file.name.lower()

    if name.endswith(".pdf"):
        return parse_pdf(file)
    elif name.endswith(".docx"):
        return parse_docx(file)
    elif name.endswith(".pptx"):
        return parse_pptx(file)
    elif name.endswith(".csv"):
        return parse_csv(file)
    elif name.endswith(".txt") or name.endswith(".md"):
        return parse_txt(file)
    else:
        return ["Unsupported file format."]

def parse_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = "\n".join([page.get_text() for page in doc])
    return split_chunks(text)

def parse_docx(file):
    doc = docx.Document(file)
    text = "\n".join([p.text for p in doc.paragraphs])
    return split_chunks(text)

def parse_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return split_chunks(text)

def parse_csv(file):
    df = pd.read_csv(file)
    text = df.to_string(index=False)
    return split_chunks(text)

def parse_txt(file):
    text = file.read().decode("utf-8")
    return split_chunks(text)

def split_chunks(text, max_length=500):
    return [text[i:i+max_length] for i in range(0, len(text), max_length)]
